#===============================================================================
#
#  Flatmap viewer and annotation tools
#
#  Copyright (c) 2019  David Brooks
#
#  Licensed under the Apache License, Version 2.0 (the "License");
#  you may not use this file except in compliance with the License.
#  You may obtain a copy of the License at
#
#      http://www.apache.org/licenses/LICENSE-2.0
#
#  Unless required by applicable law or agreed to in writing, software
#  distributed under the License is distributed on an "AS IS" BASIS,
#  WITHOUT WARRANTIES OR CONDITIONS OF ANY KIND, either express or implied.
#  See the License for the specific language governing permissions and
#  limitations under the License.
#
#===============================================================================

from math import pi as PI
import os

#===============================================================================

import pptx.shapes.connector
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.shapes import MSO_SHAPE_TYPE
from pptx.spec import autoshape_types

import svgwrite

#===============================================================================

from .extractor import GeometryExtractor, Transform
from .extractor import EMU_PER_DOT, ellipse_point
from .formula import Geometry, radians
from .presets import DML

#===============================================================================

def svg_coords(x, y):
#====================
    return (x/EMU_PER_DOT, y/EMU_PER_DOT)

def svg_units(emu):
#===================
    return emu/EMU_PER_DOT

def svg_transform(m):
#====================
    return (          m[0, 0],            m[1, 0],
                      m[0, 1],            m[1, 1],
            svg_units(m[0, 2]), svg_units(m[1, 2]))

#===============================================================================

class MakeSvgSlide(object):
    def __init__(self, slide, slide_number, slide_size, args):
        self._dwg = svgwrite.Drawing(filename=os.path.join(args.output_dir, 'slide{:02d}.svg'.format(slide_number)),
                                     size=svg_coords(slide_size[0], slide_size[1]))
        self._dwg.defs.add(self._dwg.style('.non-scaling-stroke { vector-effect: non-scaling-stroke; }'))
        self.svg_from_shapes_(slide.shapes, self._dwg)
        self._dwg.save()

    def svg_from_shapes_(self, shapes, svg_parent):
        for shape in shapes:
            if (shape.shape_type == MSO_SHAPE_TYPE.AUTO_SHAPE
             or shape.shape_type == MSO_SHAPE_TYPE.FREEFORM
             or isinstance(shape, pptx.shapes.connector.Connector)):
                self.shape_to_svg_(shape, svg_parent)

            elif shape.shape_type == MSO_SHAPE_TYPE.GROUP:
                svg_group = self._dwg.g()
                svg_group.matrix(*svg_transform(Transform(shape).matrix()))
                svg_parent.add(svg_group)
                self.svg_from_shapes_(shape.shapes, svg_group)

            elif shape.shape_type == MSO_SHAPE_TYPE.TEXT_BOX:
                pass  # or recognise name of '#layer-id' and get layer name...

            else:
                print('"{}" {} not processed...'.format(shape.name, str(shape.shape_type)))

    def shape_to_svg_(self, shape, svg_parent):
        geometry = Geometry(shape)
        for path in geometry.path_list:
            bbox = (shape.width, shape.height) if path.w is None else (path.w, path.h)
            svg_path = self._dwg.path(fill='none', stroke_width=3, class_='non-scaling-stroke') # id='sss'
            svg_path.matrix(*svg_transform(Transform(shape, bbox).matrix()))
            first_point = None
            current_point = None
            closed = False
            for c in path.getchildren():
                if   c.tag == DML('arcTo'):
                    wR = geometry.attrib_value(c, 'wR')
                    hR = geometry.attrib_value(c, 'hR')
                    stAng = radians(geometry.attrib_value(c, 'stAng'))
                    swAng = radians(geometry.attrib_value(c, 'swAng'))
                    p1 = ellipse_point(wR, hR, stAng)
                    p2 = ellipse_point(wR, hR, stAng + swAng)
                    pt = (current_point[0] - p1[0] + p2[0],
                          current_point[1] - p1[1] + p2[1])
                    large_arc_flag = 1 if swAng >= PI else 0
                    svg_path.push('A', svg_units(wR), svg_units(hR),
                                       0, large_arc_flag, 1,
                                       svg_units(pt[0]), svg_units(pt[1]))
                    current_point = pt

                elif c.tag == DML('close'):
                    if first_point is not None and current_point != first_point:
                        svg_path.push('Z')
                    closed = True
                    first_point = None
                elif c.tag == DML('cubicBezTo'):
                    coords = []
                    for p in c.getchildren():
                        pt = geometry.point(p)
                        coords.append(svg_units(pt[0]))
                        coords.append(svg_units(pt[1]))
                        current_point = pt
                    svg_path.push('C', *coords)
                elif c.tag == DML('lnTo'):
                    pt = geometry.point(c.pt)
                    svg_path.push('L', svg_units(pt[0]), svg_units(pt[1]))
                    current_point = pt
                elif c.tag == DML('moveTo'):
                    pt = geometry.point(c.pt)
                    svg_path.push('M', svg_units(pt[0]), svg_units(pt[1]))
                    if first_point is None:
                        first_point = pt
                    current_point = pt
                elif c.tag == DML('quadBezTo'):
                    coords = []
                    for p in c.getchildren():
                        pt = geometry.point(p)
                        coords.append(svg_units(pt[0]))
                        coords.append(svg_units(pt[1]))
                        current_point = pt
                    svg_path.push('Q', *coords)
                else:
                    print('Unknown path element: {}'.format(c.tag))
            if closed:
                svg_path.attribs['fill'] = '#808080'
                svg_path.attribs['opacity'] = 0.3
                svg_path.attribs['stroke'] = 'red'
            else:
                svg_path.attribs['stroke'] = 'blue'

            svg_parent.add(svg_path)

#===============================================================================

class SvgExtractor(GeometryExtractor):
    def __init__(self, args):
        super().__init__(args)
        self._slide_maker = MakeSvgSlide

#===============================================================================
