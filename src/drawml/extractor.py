#===============================================================================
#
#  Flatmap viewer and annotation tools
#
#  Copyright (c) 2019  David Brooks
#
#  Licensed under the Apache License, Version 2.0 (the "License");
#  you may not use this file except in compliance with the License.
#  You may obtain a copy of the License at
#
#      http://www.apache.org/licenses/LICENSE-2.0
#
#  Unless required by applicable law or agreed to in writing, software
#  distributed under the License is distributed on an "AS IS" BASIS,
#  WITHOUT WARRANTIES OR CONDITIONS OF ANY KIND, either express or implied.
#  See the License for the specific language governing permissions and
#  limitations under the License.
#
#===============================================================================

from math import sqrt, sin, cos, pi as PI
import os

#===============================================================================

import numpy as np

import pptx.shapes.connector
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.shapes import MSO_SHAPE_TYPE
from pptx.spec import autoshape_types


#===============================================================================

# Internal PPT units are EMUs (English Metric Units)

EMU_PER_CM  = 360000
EMU_PER_IN  = 914400

DOTS_PER_IN = 96

EMU_PER_DOT = EMU_PER_IN/DOTS_PER_IN

#===============================================================================

def cm_coords(x, y):
#===================
    return (x/EMU_PER_CM, y/EMU_PER_CM)

def ellipse_point(a, b, theta):
#==============================
    a_sin_theta = a*sin(theta)
    b_cos_theta = b*cos(theta)
    circle_radius = sqrt(a_sin_theta**2 + b_cos_theta**2)
    return (a*b_cos_theta/circle_radius, b*a_sin_theta/circle_radius)

#===============================================================================

class Transform(object):
    def __init__(self, shape, bbox=None):
        xfrm = shape.element.xfrm

        # From Section L.4.7.6 of ECMA-376 Part 1
        (Bx, By) = ((xfrm.chOff.x, xfrm.chOff.y)
                        if xfrm.chOff is not None else
                    (0, 0))
        (Dx, Dy) = ((xfrm.chExt.cx, xfrm.chExt.cy)
                        if xfrm.chExt is not None else
                    bbox)
        (Bx_, By_) = (xfrm.off.x, xfrm.off.y)
        (Dx_, Dy_) = (xfrm.ext.cx, xfrm.ext.cy)
        theta = xfrm.rot*PI/180.0
        Fx = -1 if xfrm.flipH else 1
        Fy = -1 if xfrm.flipV else 1
        T_st = np.matrix([[Dx_/Dx,      0, Bx_ - (Dx_/Dx)*Bx] if Dx != 0 else [1, 0, Bx_],
                          [     0, Dy_/Dy, By_ - (Dy_/Dy)*By] if Dy != 0 else [0, 1, By_],
                          [     0,      0,                 1]])
        U = np.matrix([[1, 0, -(Bx_ + Dx_/2.0)],
                       [0, 1, -(By_ + Dy_/2.0)],
                       [0, 0,                1]])
        R = np.matrix([[cos(theta), -sin(theta), 0],
                       [sin(theta),  cos(theta), 0],
                       [0,                    0, 1]])
        Flip = np.matrix([[Fx,  0, 0],
                          [ 0, Fy, 0],
                          [ 0,  0, 1]])
        T_rf = U.I*R*Flip*U
        self._T = T_rf*T_st

    def matrix(self):
        return self._T

#===============================================================================

class ProcessSlide(object):
    def __init__(self, slide, slide_number, args):
        self._slide = slide
        self._slide_number = slide_number
        self._args = args
        self._layer_id = 'slide{:02d}'.format(slide_number)
        self._description = 'Slide {:02d}'.format(slide_number)
        self._shape_ids = []

    @property
    def args(self):
        return self._args

    @property
    def description(self):
        return self._description

    @property
    def layer_id(self):
        return self._layer_id

    @property
    def shape_ids(self):
        return self._shape_ids

    @property
    def slide(self):
        return self._slide

    def process():
        # Override in sub-class
        pass

    def get_output(self):
        # Override in sub-class
        pass

    def save(self, filename=None):
        # Override in sub-class
        pass

    def process_shape_list(self, shapes, *args):
        for shape in shapes:
            shape.id = shape.element._nvXxPr.cNvPr.id  # FUTURE: shape.id
            # See https://github.com/scanny/python-pptx/issues/520
            shape.attribute = ''
            if shape.name.startswith('#'):
                attribs = shape.name.split()
                if len(attribs[0]) > 1:
                    shape.id = shape.name.split()[0][1:]
                    if shape.id in self._shape_ids:
                        raise KeyError('Duplicate ID {} in slide {}'
                                       .format(shape.id, self._slide_number))
                    self._shape_ids.append(shape.id)
                    if len(attribs) > 1:
                        shape.attribute = attribs[1]
            if (shape.shape_type == MSO_SHAPE_TYPE.AUTO_SHAPE
             or shape.shape_type == MSO_SHAPE_TYPE.FREEFORM
             or isinstance(shape, pptx.shapes.connector.Connector)):
                self.process_shape(shape, *args)
            elif shape.shape_type == MSO_SHAPE_TYPE.GROUP:
                self.process_group(shape, *args)
            elif shape.shape_type == MSO_SHAPE_TYPE.TEXT_BOX:
                if shape.id == 'layer-id' and shape.attribute != '':
                    self._layer_id = shape.attribute
                    if shape.text != '':
                        self._description = shape.text
            else:
                print('"{}" {} not processed...'.format(shape.name, str(shape.shape_type)))

    def geojson_from_shapes_(self, shapes, transform):
        for shape in shapes:
            if (shape.shape_type == MSO_SHAPE_TYPE.AUTO_SHAPE
             or shape.shape_type == MSO_SHAPE_TYPE.FREEFORM
             or isinstance(shape, pptx.shapes.connector.Connector)):
                self.shape_to_feature_(shape, transform)
            elif shape.shape_type == MSO_SHAPE_TYPE.GROUP:
                self.geojson_from_shapes_(shape.shapes, transform*Transform(shape).matrix())
            elif shape.shape_type == MSO_SHAPE_TYPE.TEXT_BOX:
                #print('{}: {}'.format(shape.name, shape.text)) # Recognise name of '#layer-id' and get layer name...
                pass
            else:
                print('"{}" {} not processed...'.format(shape.name, str(shape.shape_type)))

#===============================================================================

class GeometryExtractor(object):
    def __init__(self, pptx, args):
        self._ppt = Presentation(pptx)
        self._args = args
        self._slides = self._ppt.slides
        self._slide_size = [self._ppt.slide_width, self._ppt.slide_height]
        self._slide_maker = None

    def __len__(self):
        return len(self._slides)

    @property
    def slide_size(self):
        return self._slide_size

    def slide(self, slide_number):
        return self._slides[slide_number - 1]

    def slide_to_geometry(self, slide_number, save_output=True):
        slide = self.slide(slide_number)
        if self._args.debug_xml:
            xml = open(os.path.join(self._args.output_dir, 'slide{:02d}.xml'.format(slide_number)), 'w')
            xml.write(slide.element.xml)
            xml.close()
        if self._slide_maker is not None:
            slide_maker = self._slide_maker(slide, slide_number, self.slide_size, self._args)
            slide_maker.process()
            if save_output:
                slide_maker.save()
            else:
                return slide_maker

    def slides_to_geometry(self, slide_range):
        if slide_range is None:
            slide_range = range(1, len(self._slides)+1)
        elif isinstance(slide_range, int):
            slide_range = [slide_range]
        for n in slide_range:
            self.slide_to_geometry(n)

#===============================================================================
