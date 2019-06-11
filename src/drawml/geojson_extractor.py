#===============================================================================
#
#  Flatmap viewer and annotation tools
#
#  Copyright (c) 2019  David Brooks
#
#  Licensed under the Apache License, Version 2.0 (the "License");
#  you may not use this file except in compliance with the License.
#  You may obtain a copy of the License at
#
#      http://www.apache.org/licenses/LICENSE-2.0
#
#  Unless required by applicable law or agreed to in writing, software
#  distributed under the License is distributed on an "AS IS" BASIS,
#  WITHOUT WARRANTIES OR CONDITIONS OF ANY KIND, either express or implied.
#  See the License for the specific language governing permissions and
#  limitations under the License.
#
#===============================================================================

import json
import math
import os

#===============================================================================

# https://simoncozens.github.io/beziers.py/index.html
from beziers.cubicbezier import CubicBezier
from beziers.point import Point as BezierPoint
from beziers.quadraticbezier import QuadraticBezier

import pptx.shapes.connector
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.shapes import MSO_SHAPE_TYPE
from pptx.spec import autoshape_types

import numpy as np

#===============================================================================

from .extractor import GeometryExtractor, Transform
from .extractor import ellipse_point
from .formula import Geometry, radians
from .presets import DML

#===============================================================================

WORLD_PER_EMU = 0.01

def transform_point(transform, point):
    pt = transform.dot([point[0], point[1], 1.0])
    return (pt[0, 0], pt[0, 1])

def point_to_lon_lat(point):
    b = 20037508.34
    lon = point[0]
    lat = point[1]
    return (lon*180/b, math.atan(math.exp(lat*math.pi/b))*360/math.pi - 90)

def points_to_lon_lat(points):
    return [ point_to_lon_lat(pt) for pt in points ]

#===============================================================================

class MakeGeoJsonSlide(object):
    def __init__(self, slide, slide_number, slide_size, args):
        self._features = []
        self._layer_id = 'slide{:02d}'.format(slide_number)
        self._path_id = 1
        transform = np.matrix([[WORLD_PER_EMU,              0, 0],
                               [            0, -WORLD_PER_EMU, 0],
                               [            0,              0, 1]])*np.matrix([[1, 0, -slide_size[0]/2.0],
                                                                               [0, 1, -slide_size[1]/2.0],
                                                                               [0, 0,                1.0]])
        self.geojson_from_shapes_(slide.shapes, transform)
        with open(os.path.join(args.output_dir, '{}.json'.format(self._layer_id)), 'w') as output_file:
            json.dump({
                'type': 'FeatureCollection',
                'id': self._layer_id,
                'creator': 'pptx2geo',        # Add version
                'features': self._features
            }, output_file)

    def geojson_from_shapes_(self, shapes, transform):
        for shape in shapes:
            if (shape.shape_type == MSO_SHAPE_TYPE.AUTO_SHAPE
             or shape.shape_type == MSO_SHAPE_TYPE.FREEFORM
             or isinstance(shape, pptx.shapes.connector.Connector)):
                self.shape_to_feature_(shape, transform)
            elif shape.shape_type == MSO_SHAPE_TYPE.GROUP:
                self.geojson_from_shapes_(shape.shapes, transform*Transform(shape).matrix())
            elif shape.shape_type == MSO_SHAPE_TYPE.TEXT_BOX:
                #print('{}: {}'.format(shape.name, shape.text)) # Recognise name of '#layer-id' and get layer name...
                pass
            else:
                print('"{}" {} not processed...'.format(shape.name, str(shape.shape_type)))

    def shape_to_feature_(self, shape, transform):
        path_id = '{}/{}'.format(self._layer_id, self._path_id)
        self._path_id += 1
        feature = {
            'type': 'Feature',
            'id': path_id,
            'properties': {
                'id': path_id
            }
        }
        geometry = {}
        coordinates = []

        pptx_geometry = Geometry(shape)
        for path in pptx_geometry.path_list:
            bbox = (shape.width, shape.height) if path.w is None else (path.w, path.h)
            T = transform*Transform(shape, bbox).matrix()

            moved = False
            first_point = None
            current_point = None
            closed = False

            for c in path.getchildren():
                if   c.tag == DML('arcTo'):
                    wR = pptx_geometry.attrib_value(c, 'wR')
                    hR = pptx_geometry.attrib_value(c, 'hR')
                    stAng = radians(pptx_geometry.attrib_value(c, 'stAng'))
                    swAng = radians(pptx_geometry.attrib_value(c, 'swAng'))
                    p1 = ellipse_point(wR, hR, stAng)
                    p2 = ellipse_point(wR, hR, stAng + swAng)
                    pt = (current_point[0] - p1[0] + p2[0],
                          current_point[1] - p1[1] + p2[1])
                    large_arc_flag = 1 if swAng >= math.pi else 0
                    ## Arc as bezier??
                    current_point = pt

                elif c.tag == DML('close'):
                    if first_point is not None and first_point == current_point:
                        closed = True
                    first_point = None
                    # Close current pptx_geometry and start a new one...

                elif c.tag == DML('cubicBezTo'):
                    coords = [BezierPoint(*current_point)]
                    for p in c.getchildren():
                        pt = pptx_geometry.point(p)
                        coords.append(BezierPoint(*pt))
                        current_point = pt
                    bz = CubicBezier(*coords)
                    print(bz.length)
                    samples = 1000
                    coordinates.extend([transform_point(T, (pt.x, pt.y)) for pt in bz.sample(samples)])

                elif c.tag == DML('lnTo'):
                    pt = pptx_geometry.point(c.pt)
                    if moved:
                        coordinates.append(transform_point(T, current_point))
                        moved = False
                    coordinates.append(transform_point(T, pt))
                    current_point = pt

                elif c.tag == DML('moveTo'):
                    pt = pptx_geometry.point(c.pt)
                    if first_point is None:
                        first_point = pt
                    current_point = pt
                    moved = True

                elif c.tag == DML('quadBezTo'):
                    coords = [BezierPoint(*current_point)]
                    for p in c.getchildren():
                        pt = pptx_geometry.point(p)
                        coords.append(BezierPoint(*pt))
                        current_point = pt
                    bz = QuadraticBezier(*coords)
                    print(bz.length)
                    samples = 1000
                    coordinates.extend([transform_point(T, (pt.x, pt.y)) for pt in bz.sample(samples)])

                else:
                    print('Unknown path element: {}'.format(c.tag))

            lat_lon = points_to_lon_lat(coordinates)
            if closed:
                geometry['type'] = 'Polygon'
                geometry['coordinates'] = [ lat_lon ]
            else:
                geometry['type'] = 'LineString'
                geometry['coordinates'] = lat_lon

            feature['geometry'] = geometry
            self._features.append(feature)

#===============================================================================

class GeoJsonExtractor(GeometryExtractor):
    def __init__(self, args):
        super().__init__(args)
        self._slide_maker = MakeGeoJsonSlide

#===============================================================================
